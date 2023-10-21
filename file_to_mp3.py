import tkinter
from tkinter import messagebox, filedialog
import pyttsx3
from PyPDF2 import PdfReader
from pptx import Presentation
import docx
import os
import functions
import Voices



engine = pyttsx3.init()  # Initiate the pyttsx3 module

    
file_path = functions.select_file_type()
# use the select_file function from the function.py file  

file_type = file_path.split(".")[-1]  # Get the file type
    
# Get only the file name without the full path or the type
file_name = file_path.split("/")[-1].split(".")[0]

output_file_name = file_name + ".mp3"  # Create the mp3 output file path

# Get the voice id picked by the user
voice_id_picked = functions.select_voice()
# Set the voice to the one picked by the user
engine.setProperty("voice", voice_id_picked)


# Core job: File extraction and export into mp3

# PDF
if file_type == "pdf":
    text = ""

    reader = PdfReader(file_path)  # Read the select file
    count = len(reader.pages)  # Get the number of pages in the file

    for i in range(count):  # Iterate on the pages and add the text together
        page = reader.pages[i]
        page.extract_text(visitor_text=functions.visitor_body)
        text += "".join(parts)

    # Clean the text to remove any newline
    clean_text = text.strip().replace("\n", " ")

# PowerPoint
elif file_type == "pptx":
    # open pptx file
    f = open(f"{file_path}", "rb")
    prs = Presentation(f)

    text_runs = []

    # pptx reads in all text for detected shapes in powerpoint
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)

    # Join the text runs with "... " to get one text with long pauses
    text = "... ".join(text_runs)
    # Clean the text to remove any newline
    clean_text = text.strip().replace("\n", " ")

elif file_type == "docx":
    # open Word document
    doc = docx.Document(file_path)

    # read in each paragraph in file without borders
    text = [p.text for p in doc.paragraphs]
    long_text = ""
    for line in text:
        # Join the elements of the list with different conditions
        # If last letter is a period -> " "
        # If last letter is something else -> "... ""
        if len(line) > 0:
            if line[-1] != ".":
                long_text += line + "... "
            else:
                long_text += line + " "

    # Clean the text
    # Remove trailing and leading spaces
    # Replace newlines by space
    clean_text = long_text.strip().replace("\n", " ")

elif file_type == "txt":
    # Open the file and read each line
    f = open(file_path)
    text = f.readlines()
    long_text = ""
    for line in text:
        # Join the elements of the list with different conditions
        # If last letter is a period -> " "
        # If last letter is something else -> "... ""
        if len(line) > 0:
            if line[-1] != ".":
                long_text += line + "... "
            else:
                long_text += line + " "

    # Clean the text
    # Remove trailing and leading spaces
    # Replace newlines by space
    clean_text = long_text.strip().replace("\n", " ")

# ask the folder in which the user wants the .mp3 file
output_loc = functions.select_outpath() + "/" + output_file_name

engine.save_to_file(clean_text, output_file_name)  
# Save speech to file in current directory
engine.runAndWait()
engine.stop()

# move the file to the folder the user selected 
os.rename(output_file_name, output_loc)

#Tell the user that the job is done
root = tkinter.Tk()
root.withdraw()
messagebox.showinfo("Final Step","It's all done, have a good day!")
root.destroy()

# Function to use the voices folder in the exe file
def resource_path(relative_path):
    """ Get absolute path to resource, works for dev and for PyInstaller """
    try:
        # PyInstaller creates a temp folder and stores path in _MEIPASS
        base_path = sys._MEIPASS
    except Exception:
        base_path = os.path.abspath(".")

    return os.path.join(base_path, relative_path)